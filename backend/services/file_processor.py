"""
StudySage — File Processing Service
Handles PDF, DOCX, plain-text, and images (OCR).
"""
from __future__ import annotations

import io
import logging
import os
import re
from pathlib import Path

logger = logging.getLogger(__name__)

# ── Optional imports (graceful degradation if not installed) ───────────────────

try:
    import pdfplumber
    _PDFPLUMBER = True
except ImportError:
    _PDFPLUMBER = False
    logger.warning("pdfplumber not found — PDF support disabled")

try:
    import fitz  # PyMuPDF
    _PYMUPDF = True
except ImportError:
    _PYMUPDF = False

try:
    from docx import Document as DocxDocument
    _DOCX = True
except ImportError:
    _DOCX = False
    logger.warning("python-docx not found — DOCX support disabled")

try:
    from PIL import Image
    import pytesseract
    _OCR = True
except ImportError:
    _OCR = False
    logger.warning("pytesseract / Pillow not found — OCR disabled")

try:
    from pptx import Presentation
    _PPTX = True
except ImportError:
    _PPTX = False
    logger.warning("python-pptx not found — PPTX support disabled")


MAX_CHARS = int(os.getenv("MAX_EXTRACTED_CHARS", "50000"))


def _clean(text: str) -> str:
    text = re.sub(r"\s{3,}", "\n\n", text)
    text = re.sub(r"[^\S\r\n]+", " ", text)
    return text.strip()


# ── PDF ────────────────────────────────────────────────────────────────────────

def extract_pdf(data: bytes) -> tuple[str, bool]:
    """Returns (text, ocr_used)."""
    text = ""
    if _PDFPLUMBER:
        try:
            with pdfplumber.open(io.BytesIO(data)) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        text += t + "\n"
        except Exception as e:
            logger.warning("pdfplumber failed: %s", e)

    # Fallback to PyMuPDF if pdfplumber gave nothing
    if not text.strip() and _PYMUPDF:
        try:
            doc = fitz.open(stream=data, filetype="pdf")
            for page in doc:
                text += page.get_text() + "\n"
        except Exception as e:
            logger.warning("PyMuPDF failed: %s", e)

    # Last resort: OCR each page image
    if not text.strip() and _OCR and _PYMUPDF:
        logger.info("PDF has no text layer — using OCR")
        try:
            doc = fitz.open(stream=data, filetype="pdf")
            for page in doc:
                mat = fitz.Matrix(2, 2)  # 2x zoom for quality
                pix = page.get_pixmap(matrix=mat)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                text += pytesseract.image_to_string(img) + "\n"
            return _clean(text)[:MAX_CHARS], True
        except Exception as e:
            logger.warning("OCR on PDF failed: %s", e)

    return _clean(text)[:MAX_CHARS], False


# ── DOCX ───────────────────────────────────────────────────────────────────────

def extract_docx(data: bytes) -> str:
    if not _DOCX:
        return ""
    try:
        doc = DocxDocument(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return _clean("\n".join(paragraphs))[:MAX_CHARS]
    except Exception as e:
        logger.warning("DOCX extraction failed: %s", e)
        return ""


# ── PPTX ───────────────────────────────────────────────────────────────────────

def extract_pptx(data: bytes) -> str:
    if not _PPTX:
        return ""
    try:
        prs = Presentation(io.BytesIO(data))
        chunks: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    chunks.append(shape.text.strip())
        return _clean("\n".join(chunks))[:MAX_CHARS]
    except Exception as e:
        logger.warning("PPTX extraction failed: %s", e)
        return ""


# ── Image OCR ──────────────────────────────────────────────────────────────────

def extract_image(data: bytes) -> tuple[str, bool]:
    if not _OCR:
        return "", False
    try:
        img = Image.open(io.BytesIO(data))
        text = pytesseract.image_to_string(img)
        return _clean(text)[:MAX_CHARS], True
    except Exception as e:
        logger.warning("Image OCR failed: %s", e)
        return "", False


# ── Dispatcher ─────────────────────────────────────────────────────────────────

def process_file(filename: str, content_type: str, data: bytes) -> tuple[str, bool]:
    """
    Returns (extracted_text, ocr_used).
    Dispatches based on content_type / extension.
    """
    ext = Path(filename).suffix.lower()
    ct = content_type.lower()

    if ext == ".pdf" or "pdf" in ct:
        return extract_pdf(data)

    if ext in (".docx", ".doc") or "word" in ct or "docx" in ct:
        return extract_docx(data), False

    if ext in (".pptx", ".ppt") or "presentation" in ct:
        return extract_pptx(data), False

    if ext in (".txt", ".md") or ct.startswith("text/"):
        try:
            return _clean(data.decode("utf-8", errors="replace"))[:MAX_CHARS], False
        except Exception:
            return "", False

    if ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".gif") or ct.startswith("image/"):
        return extract_image(data)

    logger.warning("Unsupported file type: ext=%s ct=%s", ext, ct)
    return "", False
